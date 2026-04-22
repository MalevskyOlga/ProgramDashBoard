"""
Gantt-style PowerPoint exporter.
Draws all project tasks as a proper Gantt chart on one 16:9 slide.
"""

import os
import re
from datetime import datetime, timedelta

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _rgb(hex_color):
    h = hex_color.lstrip('#')
    return RGBColor.from_string(h)


def _parse_date(value):
    if not value:
        return None
    if isinstance(value, datetime):
        return value
    text = str(value).strip()[:19]
    for fmt in ('%Y-%m-%d', '%Y-%m-%d %H:%M:%S', '%m/%d/%Y'):
        try:
            return datetime.strptime(text, fmt)
        except ValueError:
            continue
    return None


def _sanitize(value):
    return re.sub(r'[<>:"/\\|?*\n\r]+', '_', str(value)).strip(' .') or 'project'


def _truncate(text, max_chars):
    t = str(text or '')
    return t if len(t) <= max_chars else t[:max_chars - 1] + '…'


STATUS_COLORS = {
    'Completed': 'A0AEC0',
    'In Process': '10B981',
    'Planned':    'F59E0B',
}
CRITICAL_COLOR  = 'EF4444'
MILESTONE_COLOR = '6366F1'
PHASE_BG        = '1F3A5F'
HEADER_BG       = '1B4332'


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(max(0.005, w)), Inches(max(0.005, h)))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill_hex)
    lc = line_hex or fill_hex
    s.line.color.rgb = _rgb(lc)
    return s


def _add_text(slide, text, x, y, w, h, size_pt, bold=False, italic=False, color='111827', align=PP_ALIGN.LEFT, valign_middle=True, wrap=False):
    from pptx.enum.text import MSO_AUTO_SIZE
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(max(0.05, w)), Inches(max(0.05, h)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE   # fix height, do not auto-grow
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.line_spacing = 1.0                # single line spacing — no extra inter-line gap
    p.font.name = 'Segoe UI'
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = _rgb(color)
    return tb


class GanttPptExporter:
    # Slide dimensions (inches)
    SW = 13.333
    SH = 7.5
    ML = 0.10   # margin left
    MR = 0.12   # margin right

    # Column widths (inches)
    NAME_W  = 2.55
    OWNER_W = 0.62
    GAP     = 0.05
    TL_LEFT = ML + NAME_W + OWNER_W + GAP   # 3.32
    TL_W    = SW - TL_LEFT - MR             # 9.893

    # Row heights
    HDR_H     = 0.42
    TL_HDR_H  = 0.44   # total timeline header (year row + month row)
    YR_ROW    = 0.18   # top portion: year labels
    MO_ROW    = 0.26   # bottom portion: month labels
    PHASE_H   = 0.21
    FOOT_H    = 0.20

    def export_gantt(self, project, tasks, output_folder, save_path=None,
                     title_prefix='', filename_override=None):
        try:
            # Filter to active tasks only
            tasks = [t for t in tasks if str(t.get('status', '')).strip() in ('Planned', 'In Process')]
            today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)

            # ── Organise tasks by phase ──────────────────────────────────────
            phases_order = []
            phase_map = {}
            for t in tasks:
                ph = (t.get('phase') or 'Other').strip()
                if ph not in phase_map:
                    phases_order.append(ph)
                    phase_map[ph] = []
                phase_map[ph].append(t)

            # ── Date range ──────────────────────────────────────────────────
            all_dates = []
            for t in tasks:
                for k in ('start_date', 'end_date'):
                    d = _parse_date(t.get(k))
                    if d:
                        all_dates.append(d)
            if not all_dates:
                return None

            min_d = min(all_dates)
            max_d = max(all_dates)
            proj_start = datetime(min_d.year, min_d.month, 1)
            end_m = max_d.month + 1
            end_y = max_d.year
            if end_m > 12:
                end_m = 1
                end_y += 1
            proj_end = datetime(end_y, end_m, 1)
            span_days = max(1, (proj_end - proj_start).days)

            def to_x(dt):
                if not dt:
                    return self.TL_LEFT
                frac = (dt - proj_start).days / span_days
                return self.TL_LEFT + max(0.0, min(1.0, frac)) * self.TL_W

            # ── Group by phase (gate order), sort tasks within phase by start date ──
            import re as _re
            phases_order = []
            phase_map = {}
            for t in tasks:
                ph = (t.get('phase') or 'Other').strip()
                if ph not in phase_map:
                    phases_order.append(ph)
                    phase_map[ph] = []
                phase_map[ph].append(t)

            def _gate_num(name):
                m = _re.search(r'\d+', name)
                return int(m.group()) if m else 999

            phases_order.sort(key=_gate_num)

            for ph in phases_order:
                phase_map[ph].sort(key=lambda t: _parse_date(t.get('start_date')) or proj_start)

            # ── Row height (dynamic) ────────────────────────────────────────
            n_tasks  = len(tasks)
            n_dividers = len(phases_order)
            avail_h  = self.SH - (self.HDR_H + 0.01) - self.TL_HDR_H - self.FOOT_H - 0.04 - n_dividers * self.PHASE_H
            row_h    = max(0.085, min(0.55, avail_h / max(1, n_tasks)))
            fs       = max(5.0, min(9.5, row_h * 72 * 0.38))   # font pt

            # ── Presentation ────────────────────────────────────────────────
            prs = Presentation()
            prs.slide_width  = Inches(self.SW)
            prs.slide_height = Inches(self.SH)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # ── Header band ─────────────────────────────────────────────────
            _add_rect(slide, 0, 0, self.SW, self.HDR_H, HEADER_BG)
            proj_name = title_prefix + project.get('name', 'Project Gantt')
            _add_text(slide, proj_name, 0.15, 0.05, 9.0, self.HDR_H - 0.08, 15, bold=True, color='FFFFFF')
            today_label = datetime.now().strftime('%d %b %Y')
            meta = f"PM: {project.get('manager', '')}  |  {today_label}"
            _add_text(slide, meta, 9.2, 0.05, self.SW - 9.35, self.HDR_H - 0.08, 8, color='A8D5B5', align=PP_ALIGN.RIGHT)

            # ── Legend row (inside header, bottom) ──────────────────────────
            legend = [('Completed', 'A0AEC0'), ('In Process', '10B981'),
                      ('Planned', 'F59E0B'), ('Critical', 'EF4444'), ('Milestone ◆', '6366F1')]
            lx = self.SW - self.MR
            ly = self.HDR_H - 0.19
            for lname, lcol in reversed(legend):
                tw = len(lname) * 0.065 + 0.18
                lx -= tw
                _add_rect(slide, lx, ly + 0.035, 0.09, 0.09, lcol)
                _add_text(slide, lname, lx + 0.10, ly, tw - 0.12, 0.18, 6.5, color='CCCCCC')
                lx -= 0.06

            # ── Timeline header (two rows: years on top, months below) ─────
            tl_top  = self.HDR_H + 0.01
            yr_top  = tl_top                        # year row starts here
            mo_top  = tl_top + self.YR_ROW          # month row below
            span_months = (proj_end.year - proj_start.year) * 12 + proj_end.month - proj_start.month

            # Background for entire header area
            _add_rect(slide, self.TL_LEFT, yr_top, self.TL_W, self.TL_HDR_H, 'F8FAFC')

            # ── Year row ────────────────────────────────────────────────────
            yr_alt = False
            for yr in range(proj_start.year, proj_end.year + 1):
                yr_s = datetime(yr, 1, 1)
                yr_e = datetime(yr + 1, 1, 1)
                x1 = max(to_x(proj_start), to_x(yr_s))
                x2 = min(to_x(proj_end),   to_x(yr_e))
                w  = x2 - x1
                if w <= 0:
                    yr_alt = not yr_alt
                    continue
                bg = 'DDE6F0' if yr_alt else 'EEF3F9'
                _add_rect(slide, x1, yr_top, w, self.YR_ROW, bg)
                if w > 0.25:
                    _add_text(slide, str(yr), x1 + 0.04, yr_top + 0.01,
                              w - 0.08, self.YR_ROW - 0.02,
                              7.5, bold=True, color='1E3A5F', align=PP_ALIGN.CENTER)
                # Full-height year divider line
                if yr > proj_start.year:
                    _add_rect(slide, x1, yr_top, 0.005, self.SH - yr_top - self.FOOT_H, 'A0B4C8')
                yr_alt = not yr_alt

            # ── Month / quarter row ─────────────────────────────────────────
            if span_months <= 36:
                cur = datetime(proj_start.year, proj_start.month, 1)
                alt = False
                while cur < proj_end:
                    nm = cur.month + 1
                    ny = cur.year
                    if nm > 12:
                        nm = 1
                        ny += 1
                    nx  = datetime(ny, nm, 1)
                    x   = to_x(cur)
                    w   = to_x(nx) - x
                    if alt:
                        _add_rect(slide, x, mo_top, w, self.MO_ROW, 'E4EBF5')
                    if w > 0.22:
                        label = cur.strftime('%b') if w < 0.6 else cur.strftime('%b')
                        _add_text(slide, label, x + 0.02, mo_top + 0.02,
                                  w - 0.04, self.MO_ROW - 0.04,
                                  6.5, bold=False, color='374151', align=PP_ALIGN.CENTER)
                    cur = nx
                    alt = not alt
            else:
                cur_y_yr  = proj_start.year
                cur_q_mon = ((proj_start.month - 1) // 3) * 3 + 1
                alt = False
                while True:
                    cur = datetime(cur_y_yr, cur_q_mon, 1)
                    if cur >= proj_end:
                        break
                    nq = cur_q_mon + 3
                    ny = cur_y_yr
                    if nq > 12:
                        nq -= 12
                        ny += 1
                    nx  = datetime(ny, nq, 1)
                    x   = to_x(cur)
                    w   = to_x(nx) - x
                    if alt:
                        _add_rect(slide, x, mo_top, w, self.MO_ROW, 'E4EBF5')
                    q_num = (cur_q_mon - 1) // 3 + 1
                    _add_text(slide, f'Q{q_num}', x + 0.02, mo_top + 0.02,
                              w - 0.04, self.MO_ROW - 0.04,
                              6.5, bold=False, color='374151', align=PP_ALIGN.CENTER)
                    cur_y_yr  = ny
                    cur_q_mon = nq
                    alt = not alt

            # Column header labels (Name / Owner) above tasks
            _add_rect(slide, self.ML, yr_top, self.NAME_W, self.TL_HDR_H, 'E8EDF4')
            _add_text(slide, 'Task', self.ML + 0.06, yr_top + 0.01,
                      self.NAME_W - 0.1, self.TL_HDR_H - 0.02, 7, bold=True, color='1E3A5F')
            _add_rect(slide, self.ML + self.NAME_W, yr_top, self.OWNER_W + self.GAP, self.TL_HDR_H, 'E8EDF4')
            _add_text(slide, 'Owner', self.ML + self.NAME_W + 0.04, yr_top + 0.01,
                      self.OWNER_W, self.TL_HDR_H - 0.02, 7, bold=True, color='1E3A5F')

            # Separator line below full header
            tasks_top = tl_top + self.TL_HDR_H + 0.01
            _add_rect(slide, self.ML, tl_top + self.TL_HDR_H, self.SW - self.ML - self.MR, 0.005, '6B8CAE')

            # Single vertical separator between owner column and timeline
            col_line_h = self.SH - tl_top - self.FOOT_H
            _add_rect(slide, self.TL_LEFT, tl_top, 0.004, col_line_h, '6B8CAE')

            # ── Task rows ───────────────────────────────────────────────────
            cur_y = tasks_top
            row_alt = False
            tl_right = self.TL_LEFT + self.TL_W

            def draw_divider(label, bg, fg):
                nonlocal cur_y
                _add_rect(slide, self.ML, cur_y, self.SW - self.ML - self.MR, self.PHASE_H, bg)
                _add_text(slide, label, self.ML + 0.08, cur_y + 0.01, 6.0, self.PHASE_H - 0.02, 8,
                          bold=True, color=fg)
                cur_y += self.PHASE_H

            def draw_task(task):
                nonlocal cur_y, row_alt
                if row_alt:
                    _add_rect(slide, self.ML, cur_y, self.SW - self.ML - self.MR, row_h, 'F9FAFB')
                row_alt = not row_alt

                status  = str(task.get('status', 'Planned'))
                is_crit = bool(task.get('critical'))
                is_ms   = bool(task.get('milestone'))
                t_start = _parse_date(task.get('start_date')) or proj_start
                t_end   = _parse_date(task.get('end_date'))   or t_start

                bar_color = CRITICAL_COLOR if is_crit else (MILESTONE_COLOR if is_ms else STATUS_COLORS.get(status, 'F59E0B'))

                # Textbox fills the full row height so auto_size=NONE keeps it in bounds
                ty = cur_y + 0.005
                tb_h = row_h - 0.01

                # Full task name — word-wrap within column, no truncation
                _add_text(slide, str(task.get('name', '') or ''),
                          self.ML + 0.04, ty,
                          self.NAME_W - 0.22, tb_h,
                          fs, bold=is_ms, color='111827', wrap=True)

                # Owner: first word only, hard cap 9 chars
                owner_raw = (task.get('owner') or '')
                owner_display = (owner_raw.split()[0] if owner_raw.split() else '')
                owner_display = owner_display[:9]
                _add_text(slide, owner_display,
                          self.ML + self.NAME_W + 0.04, ty,
                          self.OWNER_W - 0.10, tb_h, max(fs - 1.0, 5.0), color='6B7280')

                bx  = max(self.TL_LEFT, to_x(t_start))
                bx2 = min(tl_right, max(to_x(t_start) + 0.035, to_x(t_end)))
                bw  = max(0.035, bx2 - bx)
                bh  = row_h * 0.52
                by  = cur_y + (row_h - bh) / 2

                if is_ms:
                    dia_size = min(bh * 1.1, row_h * 0.7)
                    dia_x    = min(tl_right - dia_size, to_x(t_end) - dia_size / 2)
                    d = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND,
                                               Inches(dia_x), Inches(by - (dia_size - bh) / 2),
                                               Inches(dia_size), Inches(dia_size))
                    d.fill.solid()
                    d.fill.fore_color.rgb = _rgb(bar_color)
                    d.line.color.rgb      = _rgb(bar_color)
                else:
                    _add_rect(slide, bx, by, bw, bh, bar_color, radius=True)

                cur_y += row_h

            # Waterfall: gate groups in order, tasks chronological within each gate
            for ph in phases_order:
                if not phase_map.get(ph):
                    continue
                draw_divider(ph, '1F3A5F', 'FFFFFF')
                for task in phase_map[ph]:
                    draw_task(task)

            # ── Column cover strip between owner and timeline (clips owner overflow) ──
            cover_h = cur_y - tasks_top
            if cover_h > 0:
                owner_box_right = self.ML + self.NAME_W + 0.04 + (self.OWNER_W - 0.10)
                _add_rect(slide, owner_box_right, tasks_top,
                          self.TL_LEFT - owner_box_right + 0.004, cover_h, 'E8EDF4')

            # ── TODAY line ──────────────────────────────────────────────────
            today_x = to_x(today)
            if self.TL_LEFT <= today_x <= self.TL_LEFT + self.TL_W:
                # Line through task area only (below the full header)
                _add_rect(slide, today_x, tasks_top, 0.013,
                          self.SH - tasks_top - self.FOOT_H, 'EF4444')
                # Small tick at top of line
                _add_rect(slide, today_x - 0.01, tasks_top - 0.015, 0.033, 0.015, 'EF4444')
                # TODAY label sits in the year row — clearly above month labels
                lbl_x = max(self.TL_LEFT + 0.02, min(self.TL_LEFT + self.TL_W - 0.38, today_x - 0.19))
                _add_text(slide, 'TODAY', lbl_x, yr_top + 0.01,
                          0.38, self.YR_ROW - 0.02, 6.5, bold=True, color='EF4444', align=PP_ALIGN.CENTER)

            # ── Footer ──────────────────────────────────────────────────────
            fy = self.SH - self.FOOT_H
            _add_rect(slide, 0, fy, self.SW, self.FOOT_H, 'F3F4F6', 'E5E7EB')
            gen_ts = datetime.now().strftime('%d %b %Y %H:%M')
            _add_text(slide, f'Emerson — Program / Project Management Dashboard  |  Generated {gen_ts}',
                      0.15, fy + 0.02, self.SW - 0.3, self.FOOT_H - 0.04,
                      6.5, color='9CA3AF', align=PP_ALIGN.CENTER)

            # ── Save ────────────────────────────────────────────────────────
            timestamp = datetime.now().strftime('%d-%m-%y_%H%M')
            filename  = filename_override or f"Gantt_{_sanitize(project.get('name', 'project'))}_{timestamp}.pptx"

            if save_path:
                sp = save_path.strip().replace('\\', os.sep).replace('/', os.sep)
                if os.path.isdir(sp):
                    full_path = os.path.join(sp, filename)
                else:
                    full_path = sp if sp.lower().endswith('.pptx') else sp + '.pptx'
                parent = os.path.dirname(full_path)
                if parent:
                    os.makedirs(parent, exist_ok=True)
                prs.save(full_path)
                return full_path

            os.makedirs(output_folder, exist_ok=True)
            file_path = os.path.join(output_folder, filename)
            prs.save(file_path)
            return file_path

        except Exception as exc:
            import traceback
            print(f'Error exporting Gantt PPT: {exc}')
            traceback.print_exc()
            return None

    # ── Critical-path helpers ────────────────────────────────────────────────

    def _compute_critical_ids(self, tasks, deps, float_threshold_days=60, near_term_days=90):
        """
        Backward-pass CPM from the last gate milestone.
        Returns set of task IDs whose total float <= threshold.
        """
        task_by_id = {t['id']: t for t in tasks}

        successors   = {t['id']: [] for t in tasks}
        predecessors = {t['id']: [] for t in tasks}
        for d in deps:
            p, s = d['predecessor_id'], d['successor_id']
            if p in successors and s in predecessors:
                successors[p].append(s)
                predecessors[s].append(p)

        # Multi-source CPM: seed ALL gate milestones as anchors
        milestones = [t for t in tasks if t.get('milestone') and _parse_date(t.get('end_date'))]
        if not milestones:
            milestones = [max(tasks,
                              key=lambda t: _parse_date(t.get('end_date')) or datetime.min,
                              default=None)]
            milestones = [m for m in milestones if m]
        if not milestones:
            return set()

        # LF dict initialised from every milestone's own end_date
        lf = {m['id']: _parse_date(m['end_date']) for m in milestones}

        # Backward BFS — Bellman-Ford style (re-queue on improvement).
        # A predecessor's LF = the *start_date* of its successor (finish-to-start
        # dependency): the predecessor must be done before the successor can begin.
        # Falling back to the successor's own LF when start_date is missing.
        queue = [m['id'] for m in milestones]
        while queue:
            tid = queue.pop(0)
            task_lf = lf.get(tid)
            if not task_lf:
                continue
            # The deadline a predecessor must meet = when this task starts
            successor_start = _parse_date(task_by_id[tid].get('start_date')) if tid in task_by_id else None
            deadline_for_pred = successor_start if successor_start else task_lf
            for pid in predecessors.get(tid, []):
                if pid not in task_by_id:
                    continue
                new_lf = min(lf.get(pid, deadline_for_pred), deadline_for_pred)
                if pid not in lf or new_lf < lf[pid]:
                    lf[pid] = new_lf
                    queue.append(pid)  # re-queue on improvement

        # Condition 1: CPM float from last gate ≤ threshold
        # Condition 2: task is due within near_term_days of today (urgent regardless of gate chain)
        today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
        critical_ids = set()
        for t in tasks:
            tid   = t['id']
            t_end = _parse_date(t.get('end_date'))
            if not t_end:
                continue
            # CPM float from anchor
            if tid in lf and (lf[tid] - t_end).days <= float_threshold_days:
                critical_ids.add(tid)
            # Near-term: due within near_term_days from today — active tasks only
            active = str(t.get('status', '')).strip() in ('Planned', 'In Process')
            days_out = (t_end - today).days
            if active and -7 <= days_out <= near_term_days:
                critical_ids.add(tid)

        return critical_ids

    def export_critical_path(self, project, tasks, deps, output_folder, save_path=None):
        """
        Compute CPM and export critical-path tasks as a Gantt PPT (same layout).
        """
        # Compute which tasks are on the critical path
        critical_ids = self._compute_critical_ids(tasks, deps)
        print(f'  CPM: {len(critical_ids)} critical-path tasks out of {len(tasks)}')

        # Keep Planned / In Process tasks that are critical; always keep milestones on CP
        cp_tasks = [
            dict(t, critical=1)           # copy + force red bar
            for t in tasks
            if t['id'] in critical_ids
            and (str(t.get('status', '')).strip() in ('Planned', 'In Process')
                 or t.get('milestone'))
        ]

        if not cp_tasks:
            print('  No active critical-path tasks found — falling back to all critical-flagged tasks')
            cp_tasks = [dict(t, critical=1) for t in tasks if t.get('critical')]

        # Delegate to the same Gantt exporter with the filtered list
        # Override filename prefix
        original_name = project.get('name', 'project')
        timestamp = datetime.now().strftime('%d-%m-%y_%H%M')
        filename  = f"CriticalPath_{_sanitize(original_name)}_{timestamp}.pptx"

        result = self.export_gantt(
            project, cp_tasks, output_folder, save_path,
            title_prefix='Critical Path — ', filename_override=filename
        )
        return result
