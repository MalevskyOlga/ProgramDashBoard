"""
PPT Exporter — generates project status slides matching existing NPD meeting PPTX style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from datetime import date, datetime
import io


# ── Brand colours (matching existing PPTX) ───────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)   # Emerson dark navy
C_ACCENT     = RGBColor(0x00, 0x70, 0xC0)   # Emerson blue
C_GREEN      = RGBColor(0x70, 0xAD, 0x47)
C_ORANGE     = RGBColor(0xFF, 0x99, 0x00)
C_RED        = RGBColor(0xFF, 0x00, 0x00)
C_GREY       = RGBColor(0xD6, 0xDC, 0xE4)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_YELLOW_BG  = RGBColor(0xFF, 0xF2, 0xCC)

IMPACT_COLOR = {'High': C_RED, 'Medium': C_ORANGE, 'Low': C_GREEN}
COVERAGE_COLOR = {
    'Fully':     C_GREEN,
    'Partially': C_ORANGE,
    'No':        C_RED,
    'N/A':       C_GREY,
}
CERT_STATUS_COLOR = {
    'Approved':    C_GREEN,
    'In Progress': C_ACCENT,
    'Submitted':   C_ACCENT,
    'Failed':      C_RED,
    'Planned':     C_GREY,
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    return prs.slides.add_slide(layout)


def _txbox(slide, left, top, width, height, text='', font_size=12,
           bold=False, color=C_BLACK, bg=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = color
    if bg:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = txb.shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(sp, qn('p:spPr'))
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', f'{bg.red:02X}{bg.green:02X}{bg.blue:02X}')
    return txb


def _rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _header_band(slide, title, subtitle='', priority=None):
    """Dark header bar across the top."""
    _rect(slide, 0, 0, 10, 1.1, C_DARK_BLUE)
    pri_text = f'  #{priority}' if priority else ''
    _txbox(slide, 0.15, 0.05, 7.5, 0.55,
           text=title + pri_text, font_size=18, bold=True, color=C_WHITE)
    if subtitle:
        _txbox(slide, 0.15, 0.6, 9.5, 0.45, text=subtitle, font_size=11, color=C_GREY)


def _section_label(slide, left, top, text):
    _rect(slide, left, top, 9.5, 0.28, C_ACCENT)
    _txbox(slide, left + 0.08, top + 0.02, 9.3, 0.24,
           text=text.upper(), font_size=9, bold=True, color=C_WHITE)


def _today_str():
    return date.today().strftime('%B %d, %Y')


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def generate_project_slide(pp, tasks=None, actions=None, risks=None,
                            certs=None, updates=None, resources=None, out_path=None):
    """
    Generate a multi-slide .pptx for one project.

    Slides:
      1  Cover / Status summary
      2  Achievements & Next Steps  (card: action items)
      3  Risks register
      4  Certifications
      5  Timeline bar (gantt milestones) or Action Items (card)
      6  Resource grid
    """
    tasks     = tasks     or []
    actions   = actions   or []
    risks     = risks     or []
    certs     = certs     or []
    updates   = updates   or []
    resources = resources or []

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, pp, updates)
    _slide_achievements(prs, pp, actions if pp.get('management_type') != 'gantt' else [])
    _slide_risks(prs, pp, risks)
    _slide_certs(prs, pp, certs)
    if pp.get('management_type') == 'gantt':
        _slide_timeline(prs, pp, tasks)
    else:
        _slide_action_items(prs, pp, actions)
    _slide_resources(prs, pp, resources)

    if out_path:
        prs.save(out_path)
    else:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover / Status
# ══════════════════════════════════════════════════════════════════════════════

def _slide_cover(prs, pp, updates):
    slide = _add_slide(prs)

    name     = pp.get('name', '')
    leader   = pp.get('leader', '')
    priority = pp.get('priority', '')
    gate     = pp.get('next_gate', '')
    launch   = pp.get('planned_launch', '')
    obj      = pp.get('objective', '')
    segment  = pp.get('business_segment', '')
    proc     = pp.get('process_type', '')
    status   = pp.get('status_text', '') or ''
    nxt      = pp.get('next_steps', '') or ''

    _header_band(slide, name, f'Project Status Update  ·  {_today_str()}', priority=priority)

    # Meta pills row
    pills = [
        ('Leader', leader or '—'),
        ('Process', proc or '—'),
        ('Next Gate', gate or '—'),
        ('Launch', launch or '—'),
        ('Objective', obj or '—'),
        ('Segment', segment or '—'),
    ]
    x = 0.15
    for label, val in pills:
        _rect(slide, x, 1.2, 1.5, 0.55, C_GREY)
        _txbox(slide, x + 0.05, 1.22, 1.4, 0.22, label, font_size=7, bold=True, color=C_DARK_BLUE)
        _txbox(slide, x + 0.05, 1.44, 1.4, 0.28, val, font_size=9, color=C_BLACK)
        x += 1.6

    # Status section
    _section_label(slide, 0.15, 1.9, 'Current Status')
    _txbox(slide, 0.2, 2.25, 9.3, 1.2, text=status, font_size=11, color=C_BLACK, wrap=True)

    # Next Steps section
    _section_label(slide, 0.15, 3.55, 'Next Steps')
    _txbox(slide, 0.2, 3.9, 9.3, 1.0, text=nxt, font_size=11, color=C_BLACK, wrap=True)

    # Recent updates
    if updates:
        _section_label(slide, 0.15, 5.05, 'Recent Updates')
        y = 5.4
        for u in updates[:4]:
            ts  = u.get('created_at', '')[:10]
            auth = u.get('author', '')
            line = f'{ts}  {auth}:  {u.get("content", "")}'
            _txbox(slide, 0.2, y, 9.3, 0.3, line, font_size=9, color=C_BLACK)
            y += 0.3


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Achievements & Next Steps
# ══════════════════════════════════════════════════════════════════════════════

def _slide_achievements(prs, pp, actions):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Achievements & Next Steps  ·  {_today_str()}',
                 priority=pp.get('priority'))

    _section_label(slide, 0.15, 1.15, 'Achievements')
    # Use status_text as achievements if no dedicated field
    ach = pp.get('status_text', '') or ''
    _txbox(slide, 0.2, 1.5, 9.3, 1.5, text=ach, font_size=11, color=C_BLACK, wrap=True)

    _section_label(slide, 0.15, 3.15, 'Next Steps')
    nxt = pp.get('next_steps', '') or ''
    _txbox(slide, 0.2, 3.5, 9.3, 1.2, text=nxt, font_size=11, color=C_BLACK, wrap=True)

    # For card projects: show key open action items
    if actions:
        open_actions = [a for a in actions if a.get('status') not in ('Done',)][:6]
        if open_actions:
            _section_label(slide, 0.15, 4.85, 'Open Action Items')
            y = 5.2
            for a in open_actions:
                status_icon = '▸' if a.get('status') == 'In Progress' else '○'
                due = a.get('due_date', '') or ''
                owner = a.get('owner', '') or ''
                line = f'{status_icon}  {a.get("description", "")}   [{owner}]  Due: {due}'
                _txbox(slide, 0.25, y, 9.1, 0.28, line, font_size=9, color=C_BLACK)
                y += 0.28


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Risks
# ══════════════════════════════════════════════════════════════════════════════

def _slide_risks(prs, pp, risks):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Top Risks  ·  {_today_str()}', priority=pp.get('priority'))

    if not risks:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No risks recorded.', font_size=13, color=C_GREY)
        return

    # Table header
    col_x  = [0.15, 4.5,  6.1,  7.7,  8.9]
    col_w  = [4.2,  1.5,  1.5,  1.1,  0.9]
    headers = ['Risk Description', 'Impact', 'Probability', 'Mitigation (brief)', 'Status']
    y = 1.25
    _rect(slide, 0.15, y, 9.7, 0.3, C_DARK_BLUE)
    for i, h in enumerate(headers):
        _txbox(slide, col_x[i] + 0.05, y + 0.03, col_w[i], 0.25,
               h, font_size=8, bold=True, color=C_WHITE)
    y += 0.32

    for risk in risks[:10]:
        row_bg = C_WHITE if risks.index(risk) % 2 == 0 else RGBColor(0xF2, 0xF2, 0xF2)
        _rect(slide, 0.15, y, 9.7, 0.42, row_bg)

        # Impact dot
        imp   = risk.get('impact', 'Medium')
        prob  = risk.get('probability', 'Medium')
        mitig = (risk.get('mitigation') or '')[:60]
        status = risk.get('status', 'Open')

        _txbox(slide, col_x[0] + 0.05, y + 0.04, col_w[0], 0.36,
               risk.get('description', ''), font_size=8.5, color=C_BLACK)

        for col_i, val, in ((1, imp), (2, prob)):
            dot_col = IMPACT_COLOR.get(val, C_GREY)
            _rect(slide, col_x[col_i] + 0.05, y + 0.1, 0.12, 0.22, dot_col)
            _txbox(slide, col_x[col_i] + 0.2, y + 0.05, col_w[col_i] - 0.2, 0.35,
                   val, font_size=8.5, color=C_BLACK)

        _txbox(slide, col_x[3] + 0.05, y + 0.04, col_w[3], 0.36,
               mitig, font_size=7.5, color=C_BLACK)
        status_col = C_RED if status == 'Open' else (C_GREEN if status == 'Closed' else C_ORANGE)
        _txbox(slide, col_x[4] + 0.05, y + 0.04, col_w[4], 0.36,
               status, font_size=8, bold=True, color=status_col)
        y += 0.44


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Certifications
# ══════════════════════════════════════════════════════════════════════════════

def _slide_certs(prs, pp, certs):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Certification Status  ·  {_today_str()}', priority=pp.get('priority'))

    if not certs:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No certifications recorded.', font_size=13, color=C_GREY)
        return

    col_x = [0.15, 2.5,  4.2,  5.8,  7.4,  8.8]
    col_w = [2.2,  1.6,  1.5,  1.5,  1.3,  1.0]
    headers = ['Certification', 'Agency', 'Type', 'Expected Date', 'Actual Date', 'Status']
    y = 1.25
    _rect(slide, 0.15, y, 9.7, 0.3, C_DARK_BLUE)
    for i, h in enumerate(headers):
        _txbox(slide, col_x[i] + 0.05, y + 0.03, col_w[i], 0.25,
               h, font_size=8, bold=True, color=C_WHITE)
    y += 0.32

    for cert in certs[:12]:
        row_bg = C_WHITE if certs.index(cert) % 2 == 0 else RGBColor(0xF2, 0xF2, 0xF2)
        _rect(slide, 0.15, y, 9.7, 0.38, row_bg)

        vals = [cert.get('cert_name', ''), cert.get('agency', ''), cert.get('cert_type', ''),
                cert.get('expected_date', ''), cert.get('actual_date', ''), cert.get('status', '')]
        for i, v in enumerate(vals):
            txt_col = CERT_STATUS_COLOR.get(v, C_BLACK) if i == 5 else C_BLACK
            bold    = i == 5
            _txbox(slide, col_x[i] + 0.05, y + 0.04, col_w[i], 0.32,
                   str(v), font_size=8.5, bold=bold, color=txt_col)
        y += 0.40


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5a — Timeline bar (gantt milestones)
# ══════════════════════════════════════════════════════════════════════════════

def _slide_timeline(prs, pp, tasks):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Project Timeline  ·  {_today_str()}', priority=pp.get('priority'))

    if not tasks:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No gantt data.', font_size=13, color=C_GREY)
        return

    # Gather milestones
    milestones = [t for t in tasks if t.get('milestone')]
    phases     = sorted({t.get('phase', '') for t in tasks if t.get('phase')})

    # Date range
    all_dates = []
    for t in tasks:
        for f in ('start_date', 'end_date'):
            try:
                all_dates.append(date.fromisoformat(t[f]))
            except Exception:
                pass
    if not all_dates:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No date data in gantt.', font_size=13, color=C_GREY)
        return

    d_min  = min(all_dates)
    d_max  = max(all_dates)
    total  = max((d_max - d_min).days, 1)
    today  = date.today()

    chart_left  = 0.15
    chart_top   = 1.3
    chart_width = 9.5
    row_h       = 0.35
    label_w     = 2.5

    # Year/quarter ruler
    _draw_time_ruler(slide, d_min, d_max, chart_left + label_w, chart_top,
                     chart_width - label_w)

    y = chart_top + 0.4
    for phase in phases:
        phase_tasks = [t for t in tasks if t.get('phase') == phase and not t.get('tailed_out')]
        if not phase_tasks:
            continue
        _rect(slide, chart_left, y, label_w - 0.05, row_h * 0.85,
              C_GREY if '#' not in phase else C_GREY)
        _txbox(slide, chart_left + 0.05, y + 0.04, label_w - 0.1, row_h * 0.7,
               phase, font_size=8, bold=True, color=C_DARK_BLUE)

        for t in phase_tasks[:6]:
            try:
                ts = date.fromisoformat(t.get('start_date') or d_min.isoformat())
                te = date.fromisoformat(t.get('end_date')   or d_max.isoformat())
            except Exception:
                continue
            bx = chart_left + label_w + chart_width * (ts - d_min).days / total * (1 - label_w / chart_width)
            bw = max(chart_width * (te - ts).days / total * (1 - label_w / chart_width), 0.05)
            is_crit = t.get('critical') or t.get('milestone')
            bar_col = C_RED if is_crit else C_ACCENT
            if t.get('status', '').lower() in ('done', 'closed', 'complete'):
                bar_col = C_GREEN
            _rect(slide, bx, y + 0.05, min(bw, chart_width - (bx - chart_left - label_w)),
                  row_h * 0.6, bar_col)
        y += row_h

        if y > 7.0:
            break

    # Today line
    today_x = chart_left + label_w + chart_width * (today - d_min).days / total * (1 - label_w / chart_width)
    if chart_left + label_w <= today_x <= chart_left + chart_width:
        _rect(slide, today_x, chart_top + 0.3, 0.02, y - chart_top - 0.3, C_RED)
        _txbox(slide, today_x - 0.15, chart_top + 0.12, 0.5, 0.18,
               'TODAY', font_size=6, bold=True, color=C_RED)

    # Gate milestones
    for t in milestones:
        try:
            gd = date.fromisoformat(t.get('end_date', ''))
            gx = chart_left + label_w + chart_width * (gd - d_min).days / total * (1 - label_w / chart_width)
            _rect(slide, gx - 0.04, chart_top + 0.3, 0.08, y - chart_top - 0.3,
                  RGBColor(0xFF, 0xD7, 0x00))
            _txbox(slide, gx - 0.3, chart_top + 0.12, 0.65, 0.18,
                   t.get('name', ''), font_size=6.5, bold=True,
                   color=RGBColor(0xC0, 0x80, 0x00))
        except Exception:
            pass


def _draw_time_ruler(slide, d_min, d_max, left, top, width):
    """Draw year labels across the top."""
    total = max((d_max - d_min).days, 1)
    years = range(d_min.year, d_max.year + 2)
    for yr in years:
        yr_start = date(yr, 1, 1)
        x_offset = (yr_start - d_min).days / total * width
        if 0 <= x_offset <= width:
            _txbox(slide, left + x_offset, top, 0.8, 0.28,
                   str(yr), font_size=8, bold=True, color=C_DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5b — Action Items (card projects)
# ══════════════════════════════════════════════════════════════════════════════

def _slide_action_items(prs, pp, actions):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Action Items  ·  {_today_str()}', priority=pp.get('priority'))

    if not actions:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No action items recorded.', font_size=13, color=C_GREY)
        return

    col_x = [0.15, 3.8,  5.6,  7.0,  8.3,  9.3]
    col_w = [3.5,  1.7,  1.3,  1.2,  0.9,  0.5]
    headers = ['Description', 'Owner', 'Start', 'Due Date', 'Status', '']
    y = 1.25
    _rect(slide, 0.15, y, 9.7, 0.3, C_DARK_BLUE)
    for i, h in enumerate(headers):
        _txbox(slide, col_x[i] + 0.05, y + 0.03, col_w[i], 0.25,
               h, font_size=8, bold=True, color=C_WHITE)
    y += 0.32

    status_colors = {
        'Done':        C_GREEN,
        'In Progress': C_ACCENT,
        'Blocked':     C_RED,
        'Not Started': C_GREY,
    }

    for a in actions[:15]:
        row_bg = C_WHITE if actions.index(a) % 2 == 0 else RGBColor(0xF2, 0xF2, 0xF2)
        _rect(slide, 0.15, y, 9.7, 0.38, row_bg)
        vals = [a.get('description', ''), a.get('owner', ''),
                a.get('start_date', ''), a.get('due_date', ''), a.get('status', ''), '']
        for i, v in enumerate(vals):
            txt_col = status_colors.get(v, C_BLACK) if i == 4 else C_BLACK
            _txbox(slide, col_x[i] + 0.05, y + 0.04, col_w[i], 0.32,
                   str(v), font_size=8.5, bold=(i == 4), color=txt_col)
        # Notes tooltip
        if a.get('notes'):
            _txbox(slide, col_x[0] + 0.1, y + 0.22, col_w[0] - 0.1, 0.14,
                   f'Note: {a["notes"][:80]}', font_size=6.5,
                   color=RGBColor(0x80, 0x80, 0x80))
        y += 0.40
        if y > 7.1:
            break


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Resource Grid
# ══════════════════════════════════════════════════════════════════════════════

def _slide_resources(prs, pp, resources):
    slide = _add_slide(prs)
    _header_band(slide, pp.get('name', ''),
                 f'Resource Status  ·  {_today_str()}', priority=pp.get('priority'))

    if not resources:
        _txbox(slide, 0.5, 2.5, 9, 1, 'No resource data.', font_size=13, color=C_GREY)
        return

    # Build lookup
    res_map = {r['discipline']: r for r in resources}

    from config import DISCIPLINES

    # Two rows of discipline columns
    n     = len(DISCIPLINES)
    half  = (n + 1) // 2
    col_w = 9.6 / half

    for row_i in range(2):
        row_discs = DISCIPLINES[row_i * half: (row_i + 1) * half]
        y_header  = 1.25 + row_i * 2.5
        y_cov     = y_header + 0.32
        y_demand  = y_header + 0.72

        for ci, disc in enumerate(row_discs):
            x = 0.15 + ci * col_w
            r = res_map.get(disc, {})
            cov    = r.get('coverage', 'N/A')
            days   = r.get('demand_days', 0)
            manual = r.get('is_manual_override', 0)

            bg = COVERAGE_COLOR.get(cov, C_GREY)
            _rect(slide, x, y_header, col_w - 0.08, 0.28, C_DARK_BLUE)
            _txbox(slide, x + 0.04, y_header + 0.03, col_w - 0.12, 0.23,
                   disc, font_size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

            _rect(slide, x, y_cov, col_w - 0.08, 0.35, bg)
            _txbox(slide, x + 0.04, y_cov + 0.06, col_w - 0.12, 0.25,
                   cov, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

            days_label = f'{days:.0f}d' if days else '—'
            _txbox(slide, x + 0.04, y_demand, col_w - 0.12, 0.28,
                   days_label, font_size=9, color=C_BLACK, align=PP_ALIGN.CENTER)
            if manual:
                _txbox(slide, x + 0.04, y_demand + 0.18, col_w - 0.12, 0.18,
                       '(manual)', font_size=6, color=C_GREY, align=PP_ALIGN.CENTER)
