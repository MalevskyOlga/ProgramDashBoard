"""
Risk Register PPT Exporter
Generates a PowerPoint slide deck styled after the NPD meeting risk slide format.
Columns: Risk Statement | Probability | Impact | Mitigation/Action Plan
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pptx.oxml.ns import qn

# ── Slide dimensions (16" × 9") ─────────────────────────────────────────────
SLIDE_W = 14630400
SLIDE_H = 8229600

MARGIN_L = 490000
MARGIN_R = 490000
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R   # 13650400

# ── Column positions (absolute left, width) ──────────────────────────────────
COL_RISK_L  = 560000
COL_RISK_W  = 5080000

COL_PROB_L  = 5780000     # header label left
COL_PROB_W  = 1350000     # header label width
BADGE_PROB_L = 5850000    # badge left

COL_IMP_L   = 7280000     # header label left
COL_IMP_W   = 1050000     # header label width
BADGE_IMP_L = 7300000     # badge left

COL_MIT_L   = 8480000
COL_MIT_W   = SLIDE_W - COL_MIT_L - MARGIN_R   # ≈5660400

# ── Header bar ───────────────────────────────────────────────────────────────
HDR_T = 325783
HDR_H = 468173

# ── Row layout ───────────────────────────────────────────────────────────────
FIRST_ROW_T    = 837291
ROW_H          = 1450000   # fixed row height
ROWS_PER_SLIDE = 5

# ── Badge ────────────────────────────────────────────────────────────────────
BADGE_W = 900000
BADGE_H = 310000

# ── Colors ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x28, 0x3A, 0x5E)
TEXT_DARK   = RGBColor(0x2B, 0x2B, 0x2B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ODD     = RGBColor(0xFE, 0xF2, 0xF2)   # light pink
ROW_EVEN    = RGBColor(0xFF, 0xFB, 0xF0)   # light cream

_BADGE_MAP = {
    'very high': (RGBColor(0x99, 0x1B, 0x1B), 'VERY HIGH'),
    'high':      (RGBColor(0xDC, 0x45, 0x45), 'HIGH'),
    'medium':    (RGBColor(0xE8, 0x9B, 0x1C), 'MED'),
    'low':       (RGBColor(0x22, 0xC5, 0x5E), 'LOW'),
}
_PROB_SCORE   = {'Low': 2, 'Medium': 3, 'High': 4, 'Very High': 5}
_IMPACT_SCORE = {'Low': 2, 'Medium': 3, 'High': 4, 'Very High': 5}


def _rpn(prob, impact):
    return _PROB_SCORE.get(prob, 3) * _IMPACT_SCORE.get(impact, 3)


def _badge_info(level):
    return _BADGE_MAP.get((level or '').lower(), (RGBColor(0x94, 0xA3, 0xB8), (level or 'N/A').upper()))


def _strip_style(sp):
    """Remove python-pptx's auto-added <p:style> element that references theme colors.
    Our shapes use fully explicit fill/stroke/font, so the style element is redundant
    and can cause strict-mode validation failures in some Office builds."""
    style_el = sp._element.find(qn('p:style'))
    if style_el is not None:
        sp._element.remove(style_el)


def _rect(slide, left, top, width, height, fill_rgb):
    """Add a solid rectangle with no border."""
    sp = slide.shapes.add_shape(MSO.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_rgb
    sp.line.fill.background()
    _strip_style(sp)
    return sp


def _textbox(slide, left, top, width, height, text,
             size_pt=10, bold=False, color=TEXT_DARK,
             font='Arial', align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box supporting multi-line text (\\n split)."""
    txb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap

    lines = str(text or '').split('\n')
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


def _badge(slide, left, top, label, fill_rgb):
    """Add a rounded-rectangle badge with centered white bold text."""
    sp = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(BADGE_W), Emu(BADGE_H))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_rgb
    sp.line.fill.background()
    _strip_style(sp)

    tf = sp.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name  = 'Arial'
    run.font.size  = Pt(9)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    return sp


def _build_slide(prs, project_name, risks_batch, slide_num, total_slides):
    blank = prs.slide_layouts[6]   # truly blank layout
    slide = prs.slides.add_slide(blank)

    # ── Title ────────────────────────────────────────────────────────────────
    suffix = f' ({slide_num}/{total_slides})' if total_slides > 1 else ''
    _textbox(slide, 505888, 18000, 13167360, 560000,
             f'{project_name} — Top Risks{suffix}',
             size_pt=22, bold=True,
             color=RGBColor(0x1B, 0x2A, 0x4A), font='Arial Black')

    # ── Header bar ───────────────────────────────────────────────────────────
    _rect(slide, MARGIN_L, HDR_T, CONTENT_W, HDR_H, NAVY)

    hdr_ty = HDR_T + 70000
    hdr_th = HDR_H - 80000
    _textbox(slide, COL_RISK_L + 40000, hdr_ty, COL_RISK_W, hdr_th,
             'Risk', size_pt=12, bold=True, color=WHITE)
    _textbox(slide, COL_PROB_L, hdr_ty, COL_PROB_W, hdr_th,
             'Probability', size_pt=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, COL_IMP_L, hdr_ty, COL_IMP_W, hdr_th,
             'Impact', size_pt=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, COL_MIT_L, hdr_ty, COL_MIT_W, hdr_th,
             'Mitigation / Action Plan', size_pt=12, bold=True, color=WHITE)

    # ── Divider lines (thin navy line between columns in header) ─────────────
    for lx in (COL_PROB_L - 20000, COL_IMP_L - 20000, COL_MIT_L - 20000):
        _rect(slide, lx, HDR_T, 18000, HDR_H, RGBColor(0x1A, 0x28, 0x48))

    # ── Risk rows ─────────────────────────────────────────────────────────────
    row_top = FIRST_ROW_T
    for i, risk in enumerate(risks_batch):
        bg = ROW_ODD if i % 2 == 0 else ROW_EVEN

        # Background band
        _rect(slide, MARGIN_L, row_top, CONTENT_W, ROW_H, bg)

        # Thin column separator lines
        for lx in (COL_PROB_L - 20000, COL_IMP_L - 20000, COL_MIT_L - 20000):
            _rect(slide, lx, row_top, 12000, ROW_H, RGBColor(0xE2, 0xE8, 0xF0))

        # Risk number (small) + title
        _textbox(slide, COL_RISK_L, row_top + 60000, 180000, ROW_H - 80000,
                 f'#{i+1}', size_pt=8, color=RGBColor(0x94, 0xA3, 0xB8))
        _textbox(slide, COL_RISK_L + 200000, row_top + 50000, COL_RISK_W - 220000, ROW_H - 80000,
                 risk.get('title', ''), size_pt=9.5, color=TEXT_DARK)

        # Probability badge (vertically centered)
        badge_top = row_top + (ROW_H - BADGE_H) // 2
        prob_rgb, prob_lbl = _badge_info(risk.get('probability', 'Medium'))
        _badge(slide, BADGE_PROB_L, badge_top, prob_lbl, prob_rgb)

        # Impact badge
        imp_rgb, imp_lbl = _badge_info(risk.get('impact', 'Medium'))
        _badge(slide, BADGE_IMP_L, badge_top, imp_lbl, imp_rgb)

        # RPN score (small, below badges)
        rv = _rpn(risk.get('probability', 'Medium'), risk.get('impact', 'Medium'))
        rpn_color = (RGBColor(0xDC, 0x45, 0x45) if rv >= 16
                     else RGBColor(0xE8, 0x9B, 0x1C) if rv >= 9
                     else RGBColor(0x22, 0xC5, 0x5E))
        _textbox(slide, BADGE_PROB_L, badge_top + BADGE_H + 40000, BADGE_IMP_L + BADGE_W - BADGE_PROB_L, 200000,
                 f'RPN {rv}', size_pt=7.5, bold=True, color=rpn_color, align=PP_ALIGN.CENTER)

        # Mitigation / Action Plan text
        mit = risk.get('mitigation') or ''
        _textbox(slide, COL_MIT_L + 40000, row_top + 50000, COL_MIT_W - 60000, ROW_H - 80000,
                 mit, size_pt=9, color=TEXT_DARK)

        row_top += ROW_H

    return slide


def export_risks_to_pptx(project_name, risks):
    """Generate PPT risk slides sorted by RPN desc, open risks first.
    Returns BytesIO ready for HTTP response.
    """
    # Sort: open first, then by RPN descending
    def sort_key(r):
        closed = 1 if (r.get('status') or '').lower() == 'closed' else 0
        return (closed, -_rpn(r.get('probability', 'Medium'), r.get('impact', 'Medium')))

    risks_sorted = sorted(risks, key=sort_key)

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    # Remove the 'type' attribute from <p:sldSz> — python-pptx's default template sets
    # type="screen4x3" even when dimensions are 16:9, which causes PowerPoint to flag
    # the file as needing repair. The original NPD Meeting deck omits the type attribute
    # entirely (treated as custom/no predefined layout), so we match that behavior.
    sld_sz = prs.element.find(qn('p:sldSz'))
    if sld_sz is not None and sld_sz.get('type'):
        del sld_sz.attrib['type']

    batches = [risks_sorted[i:i + ROWS_PER_SLIDE]
               for i in range(0, max(1, len(risks_sorted)), ROWS_PER_SLIDE)]
    total = len(batches)

    for idx, batch in enumerate(batches):
        _build_slide(prs, project_name, batch, idx + 1, total)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
