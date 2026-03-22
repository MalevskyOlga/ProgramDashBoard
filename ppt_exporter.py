"""
PowerPoint exporter for the standalone project schematic schedule.
"""

from datetime import datetime
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


LANE_ORDER = [
    {"key": "discovery", "label": "Discovery & Scope", "color": "48A86D"},
    {"key": "preliminary", "label": "Architecture & Design", "color": "6D8FBD"},
    {"key": "development", "label": "Development & Verification", "color": "E0A12F"},
    {"key": "certifications", "label": "Compliance & Approvals", "color": "4AAEB7"},
    {"key": "launch", "label": "Manufacturing Readiness & Launch", "color": "405579"},
]

DISCOVERY_KEYWORDS = ["discovery", "planning", "voice of the customer", "requirements", "scope", "strategy", "roadmap"]
CERTIFICATION_KEYWORDS = ["certif", "compliance", "approval", "approvals", "csa", "regulatory", "trade compliance", "sil"]
LAUNCH_KEYWORDS = ["manufact", "mfr", "build", "testing", "test", "launch", "supplier", "sourcing", "supply", "inventory", "production", "fct", "dvt", "dv ", " dv", "ev ", " ev", "mvt", "svt", "pilot"]
LAUNCH_CALLOUT_KEYWORDS = ["manufacturing preparation", "manufacturing prep", "manufacturing", "mfr readiness", "readiness", "industrialization", "production"]
REVIEW_KEYWORDS = ["pdr", "cdr", "fdr", "trr", "mrr", "review", "architecture", "srd", "requirements"]


def _parse_date(value):
    if not value:
        return None

    if isinstance(value, datetime):
        return value

    text = str(value).strip()
    if not text:
        return None

    for fmt in ("%Y-%m-%d", "%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S", "%Y-%m-%dT%H:%M:%S.%f"):
        try:
            return datetime.strptime(text[:len(fmt)], fmt)
        except ValueError:
            continue

    try:
        return datetime.fromisoformat(text.replace("Z", "+00:00")).replace(tzinfo=None)
    except ValueError:
        return None


def _rgb(hex_color):
    return RGBColor.from_string(hex_color)


def _sanitize_filename(value):
    return re.sub(r'[<>:"/\\|?*]+', "_", value).strip(" .") or "project"


def _includes_keyword(text, keywords):
    lowered = str(text or "").lower()
    return any(keyword in lowered for keyword in keywords)


def _get_task_gate_number(task):
    for value in (task.get("name"), task.get("phase")):
        match = re.search(r"gate\s*(\d+)", str(value or ""), re.IGNORECASE)
        if match:
            return int(match.group(1))
    return None


def _is_mandatory_gate(task):
    gate_number = _get_task_gate_number(task)
    return bool(task.get("milestone")) and gate_number in (2, 3, 4, 5) and "gate" in str(task.get("name") or "").lower()


def _get_gate_timeline_date(task, gate_sign_offs):
    gate_key = task.get("phase") or task.get("name")
    sign_off = gate_sign_offs.get(gate_key) or gate_sign_offs.get(task.get("name")) or {}
    return sign_off.get("sign_off_date") or task.get("end_date") or task.get("start_date")


def _classify_lane(task):
    gate_number = _get_task_gate_number(task)
    search_text = f"{task.get('name', '')} {task.get('phase', '')}"

    if _includes_keyword(search_text, CERTIFICATION_KEYWORDS):
        return "certifications"
    if _includes_keyword(search_text, LAUNCH_KEYWORDS):
        return "launch"
    if (gate_number and gate_number <= 2) or _includes_keyword(search_text, DISCOVERY_KEYWORDS):
        return "discovery"
    if (gate_number and gate_number <= 3) or "gate 3" in str(task.get("phase") or "").lower():
        return "preliminary"
    if (gate_number and gate_number >= 4) or "gate 4" in str(task.get("phase") or "").lower() or "gate 5" in str(task.get("phase") or "").lower():
        return "development"
    return "preliminary"


def _get_segment_bounds(task_list):
    dated_tasks = [task for task in task_list if _parse_date(task.get("start_date")) or _parse_date(task.get("end_date"))]
    if not dated_tasks:
        return None

    starts = [_parse_date(task.get("start_date")) or _parse_date(task.get("end_date")) for task in dated_tasks]
    ends = [_parse_date(task.get("end_date")) or _parse_date(task.get("start_date")) for task in dated_tasks]
    return {"start": min(starts), "end": max(ends)}


def _get_mandatory_gate(mandatory_gates, gate_number):
    for gate in mandatory_gates:
        if gate["gateNumber"] == gate_number:
            return gate
    return None


def _get_highlight_date(task, lane_key, mandatory_gates):
    start = _parse_date(task.get("start_date"))
    end = _parse_date(task.get("end_date")) or start
    fallback_date = end or start
    text = f"{task.get('name', '')} {task.get('phase', '')}"

    if lane_key == "launch" and _includes_keyword(text, LAUNCH_CALLOUT_KEYWORDS):
        gate4 = _get_mandatory_gate(mandatory_gates, 4)
        gate4_date = _parse_date(gate4["dateValue"]) if gate4 else None
        if gate4_date and start and end and start <= gate4_date <= end:
            return gate4_date
        if start:
            return start

    return fallback_date


def _should_consider_highlight(task, lane_key):
    if task.get("tailed_out") or _is_mandatory_gate(task):
        return False

    text = f"{task.get('name', '')} {task.get('phase', '')}"
    if task.get("critical") or task.get("milestone"):
        return True
    if lane_key == "certifications":
        return _includes_keyword(text, CERTIFICATION_KEYWORDS)
    if lane_key == "launch":
        return _includes_keyword(text, LAUNCH_KEYWORDS)
    if lane_key == "development":
        return _includes_keyword(text, REVIEW_KEYWORDS)
    if lane_key == "preliminary":
        return _includes_keyword(text, ["architecture", "srd", "requirements", "design"])
    if lane_key == "discovery":
        return _includes_keyword(text, DISCOVERY_KEYWORDS)
    return False


def _get_highlight_priority(task, lane_key):
    text = f"{task.get('name', '')} {task.get('phase', '')}"
    score = 0
    if task.get("critical"):
        score += 60
    if task.get("milestone"):
        score += 35
    if lane_key == "launch" and _includes_keyword(text, LAUNCH_KEYWORDS):
        score += 18
    if lane_key == "launch" and _includes_keyword(text, LAUNCH_CALLOUT_KEYWORDS):
        score += 38
    if lane_key == "certifications" and _includes_keyword(text, CERTIFICATION_KEYWORDS):
        score += 20
    if lane_key == "development" and _includes_keyword(text, REVIEW_KEYWORDS):
        score += 16
    if lane_key == "preliminary" and _includes_keyword(text, ["architecture", "srd", "requirements", "design"]):
        score += 14
    if lane_key == "discovery" and _includes_keyword(text, DISCOVERY_KEYWORDS):
        score += 12

    date = _parse_date(task.get("end_date")) or _parse_date(task.get("start_date")) or datetime(1900, 1, 1)
    return score, date


def _shorten_label(text, max_length=30):
    cleaned = re.sub(r"\s*\(.*?\)\s*", " ", str(text or ""))
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if len(cleaned) <= max_length:
        return cleaned
    return f"{cleaned[:max_length - 1].rstrip()}..."


def _get_status_text(mandatory_gates):
    today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
    for gate in mandatory_gates:
        gate_date = _parse_date(gate["dateValue"])
        if gate_date and gate_date >= today:
            return f"GATE {gate['gateNumber']} IN PROGRESS"

    if mandatory_gates:
        return f"GATE {mandatory_gates[-1]['gateNumber']} COMPLETE"
    return "SCHEMATIC SCHEDULE"


def _build_model(tasks, gate_sign_offs):
    normalized_tasks = []
    for task in tasks:
        normalized_tasks.append({
            **task,
            "critical": bool(task.get("critical")),
            "milestone": bool(task.get("milestone")) or "gate" in str(task.get("name") or "").lower(),
            "tailed_out": bool(task.get("tailed_out")),
        })

    mandatory_gates = []
    for task in normalized_tasks:
        if _is_mandatory_gate(task):
            gate_number = _get_task_gate_number(task)
            gate_date = _get_gate_timeline_date(task, gate_sign_offs)
            if gate_number and _parse_date(gate_date):
                mandatory_gates.append({
                    "gateNumber": gate_number,
                    "gateName": task.get("phase") or task.get("name"),
                    "dateValue": gate_date,
                })
    mandatory_gates.sort(key=lambda gate: _parse_date(gate["dateValue"]))

    lane_map = {
        lane["key"]: {**lane, "tasks": [], "segment": None, "highlightCandidates": [], "highlight": None}
        for lane in LANE_ORDER
    }

    for task in normalized_tasks:
        if task["tailed_out"] or _is_mandatory_gate(task):
            continue
        lane_key = _classify_lane(task)
        lane_map[lane_key]["tasks"].append(task)
        if _should_consider_highlight(task, lane_key):
            lane_map[lane_key]["highlightCandidates"].append(task)

    for lane in lane_map.values():
        lane["segment"] = _get_segment_bounds(lane["tasks"])
        ranked = [
            task for task in lane["highlightCandidates"]
            if _get_highlight_date(task, lane["key"], mandatory_gates)
        ]
        ranked.sort(
            key=lambda task: (
                _get_highlight_priority(task, lane["key"])[0],
                _get_highlight_date(task, lane["key"], mandatory_gates),
            ),
            reverse=True,
        )
        lane["highlight"] = ranked[0] if ranked else None

    lanes = [lane_map[lane["key"]] for lane in LANE_ORDER]
    missing_gates = [gate_number for gate_number in (2, 3, 4, 5) if not any(gate["gateNumber"] == gate_number for gate in mandatory_gates)]
    return {
        "lanes": lanes,
        "mandatory_gates": mandatory_gates,
        "missing_gates": missing_gates,
        "status_text": _get_status_text(mandatory_gates),
    }


def _get_range(model):
    dates = []
    for lane in model["lanes"]:
        if lane["segment"]:
            dates.extend([lane["segment"]["start"], lane["segment"]["end"]])
        if lane["highlight"]:
            highlight_date = _get_highlight_date(lane["highlight"], lane["key"], model["mandatory_gates"])
            if highlight_date:
                dates.append(highlight_date)
    for gate in model["mandatory_gates"]:
        gate_date = _parse_date(gate["dateValue"])
        if gate_date:
            dates.append(gate_date)

    if not dates:
        return None

    dates.sort()
    return {
        "start": datetime(dates[0].year, 1, 1),
        "end": datetime(dates[-1].year, 12, 31),
    }


def _get_percent(date_value, date_range):
    total = (date_range["end"] - date_range["start"]).total_seconds()
    if total <= 0:
        return 0
    return ((date_value - date_range["start"]).total_seconds() / total)


class PptExporter:
    def export_schematic(self, project, tasks, gate_sign_offs, output_folder):
        try:
            model = _build_model(tasks, {row["gate_name"]: row for row in gate_sign_offs})
            date_range = _get_range(model)

            presentation = Presentation()
            presentation.slide_width = Inches(13.333)
            presentation.slide_height = Inches(7.5)
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(8.9), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = f"{project['name']} Timeline"
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.name = "Segoe UI"
            title_paragraph.font.size = Pt(22)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = _rgb("20365D")

            subtitle_box = slide.shapes.add_textbox(Inches(0.47), Inches(0.82), Inches(6.5), Inches(0.35))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = f"{project.get('manager') or 'No manager'} | Schematic Schedule | Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.font.name = "Segoe UI"
            subtitle_paragraph.font.size = Pt(9)
            subtitle_paragraph.font.color.rgb = _rgb("6B7280")

            pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(0.25), Inches(2.25), Inches(0.58))
            pill.fill.solid()
            pill.fill.fore_color.rgb = _rgb("2FAA60")
            pill.line.color.rgb = _rgb("2FAA60")
            pill_text = pill.text_frame
            pill_text.text = model["status_text"]
            pill_paragraph = pill_text.paragraphs[0]
            pill_paragraph.alignment = PP_ALIGN.CENTER
            pill_paragraph.font.name = "Segoe UI"
            pill_paragraph.font.size = Pt(10)
            pill_paragraph.font.bold = True
            pill_paragraph.font.color.rgb = _rgb("FFFFFF")

            chips = ["Timeline uses live gantt task dates and gate sign-offs"]
            chips.extend([f"Missing Gate {gate_number}" for gate_number in model["missing_gates"]])
            chip_x = 0.47
            chip_y = 1.18
            for chip_text in chips:
                chip_width = min(2.8, max(1.7, 0.075 * len(chip_text)))
                chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(chip_x), Inches(chip_y), Inches(chip_width), Inches(0.28))
                chip.fill.solid()
                if chip_text.startswith("Missing Gate"):
                    chip.fill.fore_color.rgb = _rgb("FFF6E7")
                    chip.line.color.rgb = _rgb("FFD691")
                    text_color = _rgb("9A6300")
                else:
                    chip.fill.fore_color.rgb = _rgb("F8FBFF")
                    chip.line.color.rgb = _rgb("D7E2F3")
                    text_color = _rgb("49617F")
                chip.text_frame.text = chip_text
                chip_paragraph = chip.text_frame.paragraphs[0]
                chip_paragraph.alignment = PP_ALIGN.CENTER
                chip_paragraph.font.name = "Segoe UI"
                chip_paragraph.font.size = Pt(8)
                chip_paragraph.font.bold = True
                chip_paragraph.font.color.rgb = text_color
                chip_x += chip_width + 0.12

            chart_left = 0.45
            chart_top = 1.62
            label_width = 2.35
            track_left = chart_left + label_width
            track_width = 10.15

            if not date_range:
                empty_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(10.8), Inches(0.4))
                empty_box.text_frame.text = "No dated critical tasks or gate milestones are available for this project yet."
                empty_paragraph = empty_box.text_frame.paragraphs[0]
                empty_paragraph.font.name = "Segoe UI"
                empty_paragraph.font.size = Pt(14)
                empty_paragraph.font.color.rgb = _rgb("6B7280")
            else:
                year_row_top = chart_top
                first_year = date_range["start"].year
                last_year = date_range["end"].year
                for year in range(first_year, last_year + 1):
                    year_start = datetime(year, 1, 1)
                    year_end = datetime(year, 12, 31)
                    start_pct = _get_percent(year_start, date_range)
                    end_pct = _get_percent(year_end, date_range)
                    width_pct = max(0.02, end_pct - start_pct)
                    year_box = slide.shapes.add_textbox(
                        Inches(track_left + track_width * start_pct),
                        Inches(year_row_top),
                        Inches(track_width * width_pct),
                        Inches(0.28),
                    )
                    year_box.text_frame.text = str(year)
                    paragraph = year_box.text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.name = "Segoe UI"
                    paragraph.font.size = Pt(11)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = _rgb("20365D")

                today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
                today_pct = min(1, max(0, _get_percent(today, date_range)))
                today_x = track_left + track_width * today_pct
                today_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Inches(today_x), Inches(chart_top + 0.34), Inches(0.001), Inches(4.55))
                today_line.line.color.rgb = _rgb("FF5B45")
                today_line.line.dash_style = 1
                today_line.line.width = Pt(1.5)
                today_label = slide.shapes.add_textbox(Inches(today_x - 0.22), Inches(chart_top + 0.15), Inches(0.5), Inches(0.18))
                today_label.text_frame.text = "TODAY"
                today_p = today_label.text_frame.paragraphs[0]
                today_p.alignment = PP_ALIGN.CENTER
                today_p.font.name = "Segoe UI"
                today_p.font.size = Pt(8)
                today_p.font.bold = True
                today_p.font.color.rgb = _rgb("FF5B45")

                lane_top = chart_top + 0.44
                lane_height = 0.78
                for index, lane in enumerate(model["lanes"]):
                    row_top = lane_top + index * lane_height

                    label_box = slide.shapes.add_textbox(Inches(chart_left), Inches(row_top + 0.18), Inches(label_width - 0.2), Inches(0.28))
                    label_box.text_frame.text = lane["label"]
                    label_paragraph = label_box.text_frame.paragraphs[0]
                    label_paragraph.font.name = "Segoe UI"
                    label_paragraph.font.size = Pt(11)
                    label_paragraph.font.bold = True
                    label_paragraph.font.color.rgb = _rgb("1F3254")

                    baseline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(track_left), Inches(row_top + 0.42), Inches(track_width), Inches(0.012))
                    baseline.fill.solid()
                    baseline.fill.fore_color.rgb = _rgb("E8EEF6")
                    baseline.line.color.rgb = _rgb("E8EEF6")

                    if lane["segment"]:
                        start_pct = _get_percent(lane["segment"]["start"], date_range)
                        end_pct = _get_percent(lane["segment"]["end"], date_range)
                        width_pct = max(0.015, end_pct - start_pct)
                        segment = slide.shapes.add_shape(
                            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                            Inches(track_left + track_width * start_pct),
                            Inches(row_top + 0.31),
                            Inches(track_width * width_pct),
                            Inches(0.18),
                        )
                        segment.fill.solid()
                        segment.fill.fore_color.rgb = _rgb(lane["color"])
                        segment.line.color.rgb = _rgb(lane["color"])

                    if lane["highlight"]:
                        highlight_date = _get_highlight_date(lane["highlight"], lane["key"], model["mandatory_gates"])
                        if highlight_date:
                            highlight_pct = min(0.91, max(0.09, _get_percent(highlight_date, date_range)))
                            callout = slide.shapes.add_shape(
                                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                Inches(track_left + track_width * highlight_pct - 0.74),
                                Inches(row_top - 0.02),
                                Inches(1.48),
                                Inches(0.34),
                            )
                            callout.fill.solid()
                            callout.fill.fore_color.rgb = _rgb("FFFFFF")
                            callout.line.color.rgb = _rgb("DBE5EF")
                            callout.line.width = Pt(1)
                            callout.shadow.inherit = False
                            callout.text_frame.word_wrap = True
                            callout.text_frame.text = _shorten_label(lane["highlight"]["name"], 28)
                            callout_paragraph = callout.text_frame.paragraphs[0]
                            callout_paragraph.alignment = PP_ALIGN.CENTER
                            callout_paragraph.font.name = "Segoe UI"
                            callout_paragraph.font.size = Pt(8)
                            callout_paragraph.font.bold = True
                            callout_paragraph.font.color.rgb = _rgb("1F3254")

                            date_paragraph = callout.text_frame.add_paragraph()
                            date_paragraph.text = highlight_date.strftime("%b %Y")
                            date_paragraph.alignment = PP_ALIGN.CENTER
                            date_paragraph.font.name = "Segoe UI"
                            date_paragraph.font.size = Pt(7)
                            date_paragraph.font.italic = True
                            date_paragraph.font.color.rgb = _rgb("73839A")

                            diamond = slide.shapes.add_shape(
                                MSO_AUTO_SHAPE_TYPE.DIAMOND,
                                Inches(track_left + track_width * highlight_pct - 0.06),
                                Inches(row_top + 0.31),
                                Inches(0.12),
                                Inches(0.12),
                            )
                            diamond.fill.solid()
                            diamond.fill.fore_color.rgb = _rgb("FFD458")
                            diamond.line.color.rgb = _rgb("FFD458")

                gate_row_top = lane_top + len(model["lanes"]) * lane_height + 0.06
                for gate in model["mandatory_gates"]:
                    gate_date = _parse_date(gate["dateValue"])
                    if not gate_date:
                        continue
                    gate_pct = _get_percent(gate_date, date_range)
                    gate_x = track_left + track_width * gate_pct

                    gate_label = slide.shapes.add_textbox(Inches(gate_x - 0.33), Inches(gate_row_top), Inches(0.66), Inches(0.2))
                    gate_label.text_frame.text = f"Gate {gate['gateNumber']}"
                    gate_p = gate_label.text_frame.paragraphs[0]
                    gate_p.alignment = PP_ALIGN.CENTER
                    gate_p.font.name = "Segoe UI"
                    gate_p.font.size = Pt(9)
                    gate_p.font.bold = True
                    gate_p.font.color.rgb = _rgb("29B862")

                    pole = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(gate_x - 0.01), Inches(gate_row_top + 0.18), Inches(0.02), Inches(0.25))
                    pole.fill.solid()
                    pole.fill.fore_color.rgb = _rgb("29B862")
                    pole.line.color.rgb = _rgb("29B862")

                    flag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(gate_x + 0.01), Inches(gate_row_top + 0.18), Inches(0.22), Inches(0.12))
                    flag.fill.solid()
                    flag.fill.fore_color.rgb = _rgb("29B862")
                    flag.line.color.rgb = _rgb("29B862")

                    gate_date_box = slide.shapes.add_textbox(Inches(gate_x - 0.42), Inches(gate_row_top + 0.46), Inches(0.84), Inches(0.2))
                    gate_date_box.text_frame.text = gate_date.strftime("%d %b %Y")
                    gate_date_p = gate_date_box.text_frame.paragraphs[0]
                    gate_date_p.alignment = PP_ALIGN.CENTER
                    gate_date_p.font.name = "Segoe UI"
                    gate_date_p.font.size = Pt(8)
                    gate_date_p.font.italic = True
                    gate_date_p.font.color.rgb = _rgb("666666")

            timestamp = datetime.now().strftime("%d-%m-%y_%H%M")
            filename = f"Schematic Schedule_{_sanitize_filename(project['name'])}_{timestamp}.pptx"
            os.makedirs(output_folder, exist_ok=True)
            file_path = os.path.join(output_folder, filename)
            presentation.save(file_path)
            return file_path
        except Exception as exc:
            print(f"Error exporting schematic PPT: {exc}")
            return None
